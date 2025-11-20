from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import time
from PIL import Image, ImageEnhance, ImageOps 
import io
from werkzeug.datastructures import FileStorage
import os
import traceback
import zipfile
import difflib 
import subprocess 
from pypdf import PdfReader, PdfWriter
from pdf2docx import Converter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import fitz  # PyMuPDF
from xhtml2pdf import pisa
import pandas as pd
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from waitress import serve
import sys
import platform 

# --- Setup ---
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, expose_headers=["Content-Disposition"])

# --- Utility Functions ---
def save_temp_file(file_storage):
    filename = f"temp_{int(time.time())}_{file_storage.filename}"
    filename = "".join(x for x in filename if x.isalnum() or x in "._-")
    filepath = os.path.abspath(filename)
    file_storage.save(filepath)
    return filepath

def remove_temp_file(filepath):
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
    except:
        pass

# ==============================
# SMART LIBREOFFICE DETECTION
# ==============================
def get_libreoffice_command():
    system = platform.system()
    if system == 'Windows':
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for p in possible_paths:
            if os.path.exists(p): return p
        return 'soffice'
    else:
        return 'libreoffice'

def convert_office_to_pdf_libreoffice(input_path, output_format='pdf'):
    try:
        libreoffice_cmd = get_libreoffice_command()
        cmd = [libreoffice_cmd, '--headless', '--convert-to', output_format, '--outdir', os.path.dirname(input_path), input_path]
        print(f"Running LibreOffice: {cmd}")
        process = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        if process.returncode != 0:
            print(f"LibreOffice Error Log: {process.stderr.decode()}")
            return None
            
        base_name = os.path.splitext(input_path)[0]
        output_path = f"{base_name}.{output_format}"
        
        if not os.path.exists(output_path): return None
        return output_path
    except Exception as e:
        print(f"Office Conversion Exception: {e}")
        return None

# ==============================
# 1. IMAGE PROCESSING
# ==============================
def process_img_to_pdf(file_list, target_format, is_scan_mode=False):
    try:
        images = []
        for file in file_list:
            img_data = file.read()
            img = Image.open(io.BytesIO(img_data))
            if is_scan_mode:
                img = img.convert("L") 
                sharpener = ImageEnhance.Sharpness(img); img = sharpener.enhance(2.0)
                contrast = ImageEnhance.Contrast(img); img = contrast.enhance(1.5)
                brightness = ImageEnhance.Brightness(img); img = brightness.enhance(1.1)
                img = img.convert("RGB") 
            elif img.mode in ("RGBA", "P") and target_format in ['JPG', 'PDF']:
                img = img.convert("RGB")
            images.append(img)

        if not images: return None, None, None
        byte_arr = io.BytesIO()
        mimetype, new_filename = 'application/pdf', 'converted.pdf'

        if target_format == 'PDF':
            if len(images) > 1: images[0].save(byte_arr, format='PDF', resolution=100.0, save_all=True, append_images=images[1:])
            else: images[0].save(byte_arr, format='PDF', resolution=100.0)
        elif target_format == 'JPG':
            images[0].save(byte_arr, format='JPEG'); mimetype, new_filename = 'image/jpeg', 'converted.jpg'
        elif target_format == 'PNG':
            images[0].save(byte_arr, format='PNG'); mimetype, new_filename = 'image/png', 'converted.png'

        byte_arr.seek(0)
        return byte_arr, new_filename, mimetype
    except Exception as e:
        print(f"Img Error: {e}"); return None, None, None

# ==============================
# 2. SIGN PDF
# ==============================
def process_sign_pdf(file_list):
    try:
        pdf_file = next((f for f in file_list if f.filename.lower().endswith('.pdf')), None)
        img_file = next((f for f in file_list if f.filename.lower().endswith(('.png', '.jpg', '.jpeg'))), None)
        if not pdf_file or not img_file: return None, None, None

        img_data = img_file.read()
        packet = io.BytesIO(); can = canvas.Canvas(packet, pagesize=letter)
        can.drawImage(io.BytesIO(img_data), 400, 50, width=150, height=50, mask='auto')
        can.save(); packet.seek(0)
        
        sign_pdf = PdfReader(packet); doc_pdf = PdfReader(pdf_file); writer = PdfWriter()
        for i, page in enumerate(doc_pdf.pages):
            if i == len(doc_pdf.pages) - 1: page.merge_page(sign_pdf.pages[0])
            writer.add_page(page)
        byte_arr = io.BytesIO(); writer.write(byte_arr); byte_arr.seek(0)
        return byte_arr, "signed_document.pdf", "application/pdf"
    except Exception as e: return None, None, None

# ==============================
# 3. COMPARE PDF
# ==============================
def process_compare_pdf(file_list):
    try:
        if len(file_list) < 2: return None, None, None
        def get_text(f):
            text = ""; reader = PdfReader(f)
            for page in reader.pages: text += page.extract_text() + "\n"
            return text.splitlines()
        text1 = get_text(file_list[0]); text2 = get_text(file_list[1])
        d = difflib.HtmlDiff(); html_diff = d.make_file(text1, text2, fromdesc="File 1", todesc="File 2")
        byte_arr = io.BytesIO(); pisa.CreatePDF(html_diff, dest=byte_arr); byte_arr.seek(0)
        return byte_arr, "comparison_report.pdf", "application/pdf"
    except Exception as e: return None, None, None

# ==============================
# 4. STANDARD ROUTERS
# ==============================
def process_pdf_to_jpg(file_storage):
    temp_pdf = None
    try:
        temp_pdf = save_temp_file(file_storage)
        doc = fitz.open(temp_pdf)
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_STORED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_data = pix.tobytes("jpg")
                zf.writestr(f"page_{i+1}.jpg", img_data)
        doc.close(); remove_temp_file(temp_pdf); zip_buffer.seek(0)
        return zip_buffer, "converted_pages.zip", "application/zip"
    except Exception as e:
        if temp_pdf: remove_temp_file(temp_pdf)
        return None, None, None

# --- FIX: Added missing PDF to Word function ---
def process_pdf_to_word(file_storage):
    temp_pdf = None; temp_docx = None
    try:
        temp_pdf = save_temp_file(file_storage)
        temp_docx = os.path.splitext(temp_pdf)[0] + ".docx"
        
        cv = Converter(temp_pdf)
        cv.convert(temp_docx, start=0, end=None)
        cv.close()
        
        byte_arr = io.BytesIO()
        with open(temp_docx, "rb") as f: byte_arr.write(f.read())
        byte_arr.seek(0)
        
        remove_temp_file(temp_pdf); remove_temp_file(temp_docx)
        return byte_arr, "converted.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    except Exception as e:
        print(f"PDF to Word Error: {e}")
        if temp_pdf: remove_temp_file(temp_pdf)
        if temp_docx: remove_temp_file(temp_docx)
        return None, None, None

def process_split_pdf(file_storage):
    temp_pdf = None
    try:
        temp_pdf = save_temp_file(file_storage); doc = fitz.open(temp_pdf); zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_STORED) as zf:
            for i in range(len(doc)):
                new_doc = fitz.open(); new_doc.insert_pdf(doc, from_page=i, to_page=i)
                pdf_bytes = new_doc.tobytes(); zf.writestr(f"page_{i+1}.pdf", pdf_bytes); new_doc.close()
        doc.close(); remove_temp_file(temp_pdf); zip_buffer.seek(0)
        return zip_buffer, "split_documents.zip", "application/zip"
    except Exception as e:
        if temp_pdf: remove_temp_file(temp_pdf)
        return None, None, None

def process_html_to_pdf(file_storage):
    try:
        html_content = file_storage.read().decode('utf-8'); byte_arr = io.BytesIO()
        pisa.CreatePDF(html_content, dest=byte_arr); byte_arr.seek(0)
        return byte_arr, "webpage.pdf", "application/pdf"
    except: return None, None, None

def process_word_to_pdf(file_storage):
    temp_docx = None; temp_pdf = None
    try:
        temp_docx = save_temp_file(file_storage)
        temp_pdf = convert_office_to_pdf_libreoffice(temp_docx, 'pdf')
        if not temp_pdf or not os.path.exists(temp_pdf): raise Exception("LibreOffice conversion failed.")
        byte_arr = io.BytesIO()
        with open(temp_pdf, "rb") as f: byte_arr.write(f.read())
        byte_arr.seek(0); remove_temp_file(temp_docx); remove_temp_file(temp_pdf)
        return byte_arr, "converted.pdf", "application/pdf"
    except Exception as e: 
        print(f"Word Error: {e}")
        if temp_docx: remove_temp_file(temp_docx)
        return None, None, None

def process_pdf_to_excel(file_storage):
    temp_pdf = None
    try:
        temp_pdf = save_temp_file(file_storage); byte_arr = io.BytesIO()
        with pdfplumber.open(temp_pdf) as pdf:
            all_tables = []
            for page in pdf.pages:
                tables = page.extract_tables()
                if tables:
                    for table in tables: all_tables.extend([['' if c is None else c for c in r] for r in table])
            if not all_tables: remove_temp_file(temp_pdf); return None, None, None
            pd.DataFrame(all_tables).to_excel(byte_arr, index=False, header=False)
        byte_arr.seek(0); remove_temp_file(temp_pdf)
        return byte_arr, "data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    except: return None, None, None

def process_excel_to_pdf(file_storage):
    temp_xlsx = None
    try:
        temp_xlsx = save_temp_file(file_storage); df = pd.read_excel(temp_xlsx)
        html_str = f"<html><head><style>table{{width:100%;border-collapse:collapse}}td,th{{border:1px solid #ddd;padding:8px}}</style></head><body>{df.to_html(index=False, na_rep='')}</body></html>"
        byte_arr = io.BytesIO(); pisa.CreatePDF(html_str, dest=byte_arr)
        remove_temp_file(temp_xlsx); byte_arr.seek(0)
        return byte_arr, "sheet.pdf", "application/pdf"
    except: return None, None, None

def process_powerpoint_to_pdf(file_storage):
    temp_ppt = None; temp_pdf = None
    try:
        temp_ppt = save_temp_file(file_storage)
        temp_pdf = convert_office_to_pdf_libreoffice(temp_ppt, 'pdf')
        if not temp_pdf or not os.path.exists(temp_pdf): raise Exception("LibreOffice conversion failed.")
        byte_arr = io.BytesIO()
        with open(temp_pdf, "rb") as f: byte_arr.write(f.read())
        byte_arr.seek(0); remove_temp_file(temp_ppt); remove_temp_file(temp_pdf)
        return byte_arr, "slides.pdf", "application/pdf"
    except Exception as e: 
        print(f"PPT Error: {e}")
        if temp_ppt: remove_temp_file(temp_ppt)
        return None, None, None

def process_pdf_to_powerpoint(file_storage):
    temp_pdf = None; temp_ppt = "slides.pptx"
    try:
        temp_pdf = save_temp_file(file_storage); doc = fitz.open(temp_pdf)
        prs = Presentation(); layout = prs.slide_layouts[6]
        for i in range(len(doc)):
            page = doc.load_page(i); pix = page.get_pixmap(dpi=150)
            img_path = f"slide_{i}.png"; pix.save(img_path)
            slide = prs.slides.add_slide(layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), height=Inches(7.5))
            os.remove(img_path)
        doc.close(); prs.save(temp_ppt); byte_arr = io.BytesIO()
        with open(temp_ppt, "rb") as f: byte_arr.write(f.read())
        byte_arr.seek(0); remove_temp_file(temp_pdf); os.remove(temp_ppt)
        return byte_arr, "slides.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    except: return None, None, None

def create_watermark_pdf(text):
    packet = io.BytesIO(); can = canvas.Canvas(packet, pagesize=letter); can.setFont("Helvetica-Bold", 50)
    can.setFillColorRGB(0.5,0.5,0.5); can.setFillAlpha(0.5); can.saveState()
    can.translate(300,400); can.rotate(45); can.drawCentredString(0,0,text); can.restoreState(); can.save()
    packet.seek(0); return PdfReader(packet)

def create_pagenum_pdf(n, total):
    packet = io.BytesIO(); can = canvas.Canvas(packet, pagesize=letter); can.setFont("Helvetica", 10)
    can.drawCentredString(300, 20, f"Page {n} of {total}"); can.save(); packet.seek(0); return PdfReader(packet)

def process_pdf_tools(file_list, tool_type, custom_text=""):
    try:
        writer = PdfWriter(); byte_arr = io.BytesIO(); new_filename = "processed.pdf"
        if tool_type == 'merge-pdf':
            for f in file_list:
                f.seek(0); r = PdfReader(f)
                if r.is_encrypted: 
                    try: r.decrypt("") 
                    except: pass
                writer.append_pages_from_reader(r)
            new_filename = "merged.pdf"
        else:
            reader = PdfReader(file_list[0])
            if tool_type == 'organize-pdf':
                total = len(reader.pages); keep = []
                parts = custom_text.split(',') if custom_text else []
                if not parts: keep = list(range(total))
                else:
                    for part in parts:
                        p = part.strip()
                        if '-' in p: s,e = map(int, p.split('-')); keep.extend(range(s-1, e))
                        elif p.isdigit(): keep.append(int(p)-1)
                for i in keep: 
                    if 0 <= i < total: writer.add_page(reader.pages[i])
                new_filename = "organized.pdf"
            elif tool_type == 'crop-pdf':
                for p in reader.pages:
                    p.cropbox.lower_left = (p.cropbox.lower_left[0]+20, p.cropbox.lower_left[1]+20)
                    p.cropbox.upper_right = (p.cropbox.upper_right[0]-20, p.cropbox.upper_right[1]-20)
                    writer.add_page(p)
                new_filename = "cropped.pdf"
            elif tool_type == 'edit-pdf':
                for p in reader.pages:
                    if "/Annots" in p: del p["/Annots"]
                    writer.add_page(p)
                new_filename = "flattened_edited.pdf"
            elif tool_type == 'rotate-pdf': 
                for p in reader.pages: p.rotate(90); writer.add_page(p)
            elif tool_type == 'compress-pdf':
                for p in reader.pages: writer.add_page(p)
                for p in writer.pages: p.compress_content_streams()
                writer.add_metadata({})
            elif tool_type == 'protect-pdf': 
                writer.append_pages_from_reader(reader); writer.encrypt(custom_text if custom_text else "pdfiy")
            elif tool_type == 'unlock-pdf':
                if reader.is_encrypted: 
                    try: reader.decrypt(custom_text); reader.decrypt("") 
                    except: pass
                writer.append_pages_from_reader(reader)
            elif tool_type == 'watermark-pdf':
                wm = create_watermark_pdf(custom_text if custom_text else "pdfiy").pages[0]
                for p in reader.pages: p.merge_page(wm); writer.add_page(p)
            elif tool_type == 'page-numbers':
                total = len(reader.pages)
                for i,p in enumerate(reader.pages): p.merge_page(create_pagenum_pdf(i+1, total).pages[0]); writer.add_page(p)
            elif tool_type == 'repair-pdf': writer.append_pages_from_reader(reader)

        writer.write(byte_arr); byte_arr.seek(0)
        return byte_arr, new_filename, 'application/pdf'
    except: return None, None, None

# ==============================
# API ROUTER
# ==============================
@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files: return jsonify({"success": False, "message": "No file."}), 400
    
    uploaded_files = request.files.getlist('file')
    tool_type = request.form.get('toolType', 'img-to-pdf')
    target_format = request.form.get('targetFormat', 'PDF').upper()
    custom_text = request.form.get('customText', '').strip()

    if not uploaded_files: return jsonify({"success": False, "message": "Empty."}), 400
    print(f"Tool: {tool_type} | Files: {len(uploaded_files)}")

    # 1. SPECIAL ROUTERS
    if tool_type == 'pdf-to-word': data, name, mime = process_pdf_to_word(uploaded_files[0]) # Now defined
    elif tool_type == 'word-to-pdf': data, name, mime = process_word_to_pdf(uploaded_files[0])
    elif tool_type == 'pdf-to-jpg': data, name, mime = process_pdf_to_jpg(uploaded_files[0])
    elif tool_type == 'split-pdf': data, name, mime = process_split_pdf(uploaded_files[0])
    elif tool_type == 'html-to-pdf': data, name, mime = process_html_to_pdf(uploaded_files[0])
    elif tool_type == 'pdf-to-excel': data, name, mime = process_pdf_to_excel(uploaded_files[0])
    elif tool_type == 'excel-to-pdf': data, name, mime = process_excel_to_pdf(uploaded_files[0])
    elif tool_type == 'powerpoint-to-pdf': data, name, mime = process_powerpoint_to_pdf(uploaded_files[0])
    elif tool_type == 'pdf-to-powerpoint': data, name, mime = process_pdf_to_powerpoint(uploaded_files[0])
    elif tool_type == 'sign-pdf': data, name, mime = process_sign_pdf(uploaded_files)
    elif tool_type == 'compare-pdf': data, name, mime = process_compare_pdf(uploaded_files)
    elif tool_type == 'scan-to-pdf': data, name, mime = process_img_to_pdf(uploaded_files, 'PDF', is_scan_mode=True)

    # 2. IMAGE TOOLS
    elif tool_type in ['img-to-pdf', 'jpg-to-pdf']: 
        data, name, mime = process_img_to_pdf(uploaded_files, target_format)

    # 3. STANDARD TOOLS
    elif tool_type in ['merge-pdf', 'rotate-pdf', 'compress-pdf', 'protect-pdf', 'unlock-pdf', 'watermark-pdf', 'page-numbers', 'repair-pdf', 'organize-pdf', 'crop-pdf', 'edit-pdf']:
        data, name, mime = process_pdf_tools(uploaded_files, tool_type, custom_text)

    else:
        time.sleep(1)
        return jsonify({ "success": True, "message": f"Simulation: '{tool_type}' coming soon!" })

    if data: return send_file(data, mimetype=mime, as_attachment=True, download_name=name)
    else: return jsonify({"success": False, "message": "Processing failed or input invalid."}), 500

if __name__ == '__main__':
    # PRODUCTION SERVER
    print("Starting pdfiy production server on http://0.0.0.0:5000")
    serve(app, host='0.0.0.0', port=5000, threads=4)